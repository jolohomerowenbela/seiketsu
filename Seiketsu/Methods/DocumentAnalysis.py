from PyQt5.QtCore import *
from PyQt5.QtGui import *
from PyQt5.QtWidgets import *
from docx import Document
from pptx import Presentation
import Seiketsu.Methods.Patterns
import pandas as pd
import PyPDF2
import mimetypes as mt

class DocumentAnalysis():
    def scan(self, path):
        mimetypes = mt.guess_type(path)
        mimetypes = [mime for mime in mimetypes if mime is not None]
        for mime in mimetypes:
            if "application/pdf" in mime:
                return self.analyze_pdf(path)
            elif "vnd.openxmlformats-officedocument.wordprocessingml.document" in mime:
                return self.analyze_docx(path)
            elif "vnd.openxmlformats-officedocument.presentationml.presentation" in mime:
                return self.analyze_presentation(path)
            if "vnd.openxmlformats-officedocument.spreadsheetml.sheet" in mime:
                return self.analyze_spreadsheet(path)
            else:
                return False
    
    def analyze_pdf(self, path):
        pdfReader = PyPDF2.PdfReader(path)
        for page in pdfReader.pages:
            for word in Seiketsu.Methods.Patterns.keywords():
                if word in page.extract_text():
                    return True, word
    
    def analyze_docx(self, path):
        doc = Document(path)
        for paragraph in doc.paragraphs:
            for word in Seiketsu.Methods.Patterns.keywords():
                if word in paragraph.text:
                    return True, word
    
    def analyze_spreadsheet(self, path):
        excel = pd.read_excel(path)
        contents = excel.to_string()
        for word in Seiketsu.Methods.Patterns.keywords():
            if word in contents:
                return True, word
    
    def analyze_presentation(self, path):
        ppt = Presentation(path)
        for slide in ppt.slides:
            for shape in slide.shapes:
                for word in Seiketsu.Methods.Patterns.keywords():
                    if hasattr(shape, "text"):
                        if word in shape.text:
                            return True, word